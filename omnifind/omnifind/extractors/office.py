"""Microsoft Office 抽取器:docx / xlsx / pptx。"""
from __future__ import annotations
from pathlib import Path
from omnifind.extractors import BaseExtractor, ExtractResult, register


@register
class DocxExtractor(BaseExtractor):
    EXTENSIONS = [".docx"]
    PRIORITY = 10

    def extract(self, path: Path) -> ExtractResult:
        import docx
        import io
        # 读入内存后关闭磁盘句柄(python-docx 1.x 的 Document 不支持 with 上下文管理器),
        # 避免大批量索引时文件句柄累积泄漏
        with open(path, "rb") as fh:
            buf = io.BytesIO(fh.read())
        d = docx.Document(buf)
        paras = [p.text for p in d.paragraphs if p.text.strip()]
        # 表格文本
        for tbl in d.tables:
            for row in tbl.rows:
                cells = [c.text for c in row.cells if c.text.strip()]
                if cells:
                    paras.append("\t".join(cells))
        return ExtractResult(text="\n".join(paras), title=path.stem, metadata={"ext": ".docx"})


@register
class XlsxExtractor(BaseExtractor):
    EXTENSIONS = [".xlsx"]
    PRIORITY = 10

    def extract(self, path: Path) -> ExtractResult:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append("\t".join(cells))
        wb.close()
        return ExtractResult(text="\n".join(parts), title=path.stem, metadata={"ext": ".xlsx"})


@register
class PptxExtractor(BaseExtractor):
    EXTENSIONS = [".pptx"]
    PRIORITY = 10

    def extract(self, path: Path) -> ExtractResult:
        from pptx import Presentation
        import io
        # 读入内存后关闭磁盘句柄(python-pptx 1.x 的 Presentation 不支持 with 上下文管理器)
        with open(path, "rb") as fh:
            buf = io.BytesIO(fh.read())
        prs = Presentation(buf)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in p.runs)
                        if t.strip():
                            parts.append(t)
        return ExtractResult(text="\n".join(parts), title=path.stem, metadata={"ext": ".pptx"})


@register
class XlsExtractor(BaseExtractor):
    EXTENSIONS = [".xls"]
    PRIORITY = 10

    def extract(self, path: Path) -> ExtractResult:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        try:
            parts = []
            for sheet in wb.sheets():
                parts.append(f"=== {sheet.name} ===")
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    row_text = [str(c) if c is not None else "" for c in row]
                    if any(cell.strip() for cell in row_text if cell):
                        parts.append("\t".join(row_text))
            return ExtractResult(text="\n".join(parts), title=path.stem, metadata={"ext": ".xls"})
        finally:
            wb.release_resources()
