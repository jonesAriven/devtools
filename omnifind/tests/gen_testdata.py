"""生成多类型测试文档,用于端到端验证 L1+L2。"""
import os
from pathlib import Path

BASE = Path("/tmp/omnifind_testdata")
BASE.mkdir(parents=True, exist_ok=True)

# 1. 纯文本 + markdown + 代码
(BASE / "项目预算报告.txt").write_text(
    "2025年度项目预算总额为一千二百万元,其中研发投入占比百分之六十。", encoding="utf-8")
(BASE / "会议纪要.md").write_text(
    "# 周会纪要\n\n讨论了微服务架构拆分方案,决定采用领域驱动设计。", encoding="utf-8")
(BASE / "deploy.py").write_text(
    "# 部署脚本\ndef deploy():\n    print('starting kubernetes rollout')\n", encoding="utf-8")

# 2. PDF
import fitz
doc = fitz.open()
page = doc.new_page()
page.insert_text((72, 72), "OmniFind PDF test: quarterly financial summary 季度财务摘要")
doc.save(str(BASE / "季度财务.pdf"))
doc.close()

# 3. Word docx
import docx
d = docx.Document()
d.add_heading("产品需求文档", 0)
d.add_paragraph("本产品支持离线全文检索与语义搜索,面向个人电脑用户。")
tbl = d.add_table(rows=1, cols=2)
tbl.rows[0].cells[0].text = "功能"
tbl.rows[0].cells[1].text = "向量检索"
d.save(str(BASE / "需求文档.docx"))

# 4. Excel xlsx
import openpyxl
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "销售数据"
ws.append(["月份", "销售额"])
ws.append(["一月", 88000])
ws.append(["二月", 92000])
wb.save(str(BASE / "销售表.xlsx"))

# 5. PPTX
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "融资路演"
slide.placeholders[1].text = "我们的目标市场是中小企业知识管理"
prs.save(str(BASE / "路演PPT.pptx"))

print("测试数据生成完成:")
for f in sorted(BASE.iterdir()):
    print(f"  {f.name}  ({f.stat().st_size} bytes)")
